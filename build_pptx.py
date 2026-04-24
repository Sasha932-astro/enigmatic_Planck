"""
Assemble a PowerPoint summarizing all scenarios.

Layout per scenario (one slide):
  - Title (scenario name, cleaned)
  - Description text
  - Regression plot (main result, large, right half)
  - Diagnostic plot (left half, smaller)
  - Animation still (bottom left) with a note pointing to the GIF file

Plus a title slide and a final comparison slide.
"""
from __future__ import annotations
import pathlib
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from run_all import SCENARIOS, OUT


ROOT = pathlib.Path(__file__).resolve().parent
OUT_PPTX = ROOT / "enigmatic_Planck_scenarios.pptx"

# Widescreen 16:9
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=(0, 0, 0)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
    return tb


def fit_image(slide, path: pathlib.Path, left, top, box_w, box_h):
    """Add an image inside a (left, top, box_w, box_h) box, preserving aspect ratio."""
    with Image.open(path) as im:
        pw, ph = im.size
    box_ratio = box_w / box_h
    img_ratio = pw / ph
    if img_ratio > box_ratio:
        w = box_w
        h = int(box_w / img_ratio)
    else:
        h = box_h
        w = int(box_h * img_ratio)
    # center within box
    x = left + (box_w - w) // 2
    y = top  + (box_h - h) // 2
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def extract_midframe(gif_path: pathlib.Path) -> pathlib.Path:
    """Grab the middle frame of the animation to embed as a still."""
    png_path = gif_path.with_suffix(".midframe.png")
    with Image.open(gif_path) as im:
        n = getattr(im, "n_frames", 1)
        im.seek(n // 2)
        im.convert("RGB").save(png_path)
    return png_path


# --- Title slide --------------------------------------------------------
slide = prs.slides.add_slide(blank)
add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.2),
            "Solar Irradiance Variability Due To Rotation",
            size=36, bold=True)
add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.0),
            "Per-scenario regression of SSI on TSI across nine magnetic-feature distributions",
            size=20, color=(60, 60, 60))
add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
            "Model: solar rotation (P = 27 d, B₀ = 0). Grid: 180 × 360. n_steps = 120.",
            size=12, color=(90, 90, 90))


# --- One slide per scenario --------------------------------------------
for i, sc in enumerate(SCENARIOS, start=1):
    folder = OUT / sc["name"]
    slide = prs.slides.add_slide(blank)

    title_txt = f"{i}. {sc['name'].split('_', 1)[1].replace('_', ' ').title()}"
    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6),
                title_txt, size=26, bold=True)
    add_textbox(slide, Inches(0.4), Inches(0.8), Inches(12.5), Inches(0.55),
                sc["description"], size=13, color=(70, 70, 70))

    # Regression (MAIN result) — right half, large
    fit_image(slide, folder / "regression_big.png",
              left=Inches(6.9), top=Inches(1.5),
              box_w=Inches(6.2), box_h=Inches(5.4))
    add_textbox(slide, Inches(6.9), Inches(6.95), Inches(6.2), Inches(0.3),
                "Main result: regression slope a(λ) and correlation r(λ)",
                size=11, bold=True, color=(140, 0, 0))

    # Diagnostic — left, top
    fit_image(slide, folder / "diagnostic.png",
              left=Inches(0.3), top=Inches(1.5),
              box_w=Inches(6.4), box_h=Inches(3.8))
    add_textbox(slide, Inches(0.3), Inches(5.3), Inches(6.4), Inches(0.3),
                "Diagnostic: TSI, a(λ), r(λ), per-λ scatter with fit",
                size=10, color=(100, 100, 100))

    # Animation still — left, bottom
    still_png = extract_midframe(folder / "animation.gif")
    fit_image(slide, still_png,
              left=Inches(0.3), top=Inches(5.7),
              box_w=Inches(3.0), box_h=Inches(1.6))
    add_textbox(slide, Inches(3.4), Inches(5.7), Inches(3.3), Inches(1.6),
                f"Animation: {sc['name']}/animation.gif\n"
                f"(middle frame shown; open the .gif in a browser or Quick Look to play)",
                size=10, color=(100, 100, 100))


# --- Comparison overlay slide (full wavelength range) ------------------
slide = prs.slides.add_slide(blank)
add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6),
            "Comparison across all scenarios",
            size=26, bold=True)
add_textbox(slide, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.5),
            "Regression slope a(λ) and correlation r(λ) overlaid for all 9 scenarios (200–2000 nm).",
            size=13, color=(70, 70, 70))
fit_image(slide, OUT / "_comparison_all.png",
          left=Inches(0.4), top=Inches(1.5),
          box_w=Inches(12.5), box_h=Inches(5.6))
add_textbox(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
            "Takeaway: pure-spot scenarios flip a(λ) sign in the mid-UV; faculae-dominated and mixed-feature scenarios keep a(λ) > 0 across the visible.",
            size=11, color=(140, 0, 0))


# --- Zoomed comparison slide (350-1300 nm, easier to read) -------------
zoom_path = OUT / "_comparison_350_1300.png"
if zoom_path.exists():
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6),
                "Comparison across scenarios — zoomed (350–1300 nm)",
                size=26, bold=True)
    add_textbox(slide, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.5),
                "Same overlay restricted to 350–1300 nm for readability; the UV noise compressed the full-range scale.",
                size=13, color=(70, 70, 70))
    fit_image(slide, zoom_path,
              left=Inches(0.4), top=Inches(1.5),
              box_w=Inches(12.5), box_h=Inches(5.6))
    add_textbox(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
                "Spot-only scenarios have a(λ) near or below zero; faculae and mixed scenarios have a(λ) ≈ 1–2 peaking in the blue.",
                size=11, color=(140, 0, 0))


prs.save(str(OUT_PPTX))
print(f"Saved -> {OUT_PPTX}")
print(f"Slides: {len(prs.slides)}")
